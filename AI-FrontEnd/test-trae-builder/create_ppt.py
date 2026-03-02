from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建PPT对象
prs = Presentation()

# 标题页
slide_layout = prs.slide_layouts[0]  # 标题幻灯片
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "DeepSeek-R1: Incentivizing Reasoning Capability in LLMs via Reinforcement Learning"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle.text = "DeepSeek-AI Research Team\n2025"
subtitle.text_frame.paragraphs[0].font.size = Pt(18)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 目录页
slide_layout = prs.slide_layouts[1]  # 标题和内容
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title = slide.shapes.title
title.text = "目录"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# 添加目录内容
content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

# 添加项目符号
points = [
    "1. 研究背景与动机",
    "2. DeepSeek-R1 模型架构",
    "3. 强化学习激励机制",
    "4. 实验结果与分析",
    "5. 评估基准测试",
    "6. 结论与未来工作"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(51, 51, 51)

# 研究背景与动机
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "研究背景与动机"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    "- 大型语言模型(LLMs)在各种任务中表现出色，但推理能力仍有提升空间",
    "- 传统的监督微调方法在复杂推理任务上效果有限",
    "- 强化学习(RL)被证明可以有效提升模型的特定能力",
    "- 本研究旨在通过强化学习激励机制增强LLMs的推理能力"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)

# DeepSeek-R1 模型架构
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "DeepSeek-R1 模型架构"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    "- 基于Transformer架构的大型语言模型",
    "- 采用强化学习进行推理能力增强",
    "- 结合了多种推理技术：",
    "  • 思维链(Chain-of-Thought)",
    "  • 自我纠正(Self-Correction)",
    "  • 多步推理(Multi-step Reasoning)"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    if "•" in point:
        p.level = 1
    else:
        p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)

# 强化学习激励机制
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "强化学习激励机制"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    "- 设计了专门的奖励机制来激励模型的推理能力",
    "- 奖励信号包括：",
    "  • 推理步骤的正确性",
    "  • 最终答案的准确性",
    "  • 推理过程的完整性和逻辑性",
    "- 采用PPO(Proximal Policy Optimization)算法进行训练",
    "- 结合人类反馈强化学习(RLHF)技术"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    if "•" in point:
        p.level = 1
    else:
        p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)

# 实验结果与分析 - 图表1
slide_layout = prs.slide_layouts[5]  # 仅标题
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "实验结果与分析"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# 添加图表1
img_path = "extracted_images/page16_img1.png"
if os.path.exists(img_path):
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5.5)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

# 实验结果与分析 - 图表2
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "更多实验结果"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# 添加图表2
img_path = "extracted_images/page43_img1.png"
if os.path.exists(img_path):
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5.5)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

# 实验结果与分析 - 图表3
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "模型性能对比"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# 添加图表3
img_path = "extracted_images/page44_img1.png"
if os.path.exists(img_path):
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5.5)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)

# 评估基准测试
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "评估基准测试"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    "- 使用多种评估基准测试模型的推理能力：",
    "  • C-EVAL：覆盖52个学科的中文知识评估",
    "  • SimpleQA：事实性问答能力评估",
    "  • C-SimpleQA：中文事实性问答评估",
    "  • 其他专业领域测试（医学、法律等）",
    "- 测试结果显示模型在各个领域都表现出了强大的推理能力"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    if "•" in point:
        p.level = 1
    else:
        p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)

# 结论与未来工作
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "结论与未来工作"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    "- 成功设计了DeepSeek-R1模型，通过强化学习增强了LLMs的推理能力",
    "- 提出了有效的奖励机制，激励模型生成高质量的推理过程",
    "- 在多种评估基准上取得了优异的成绩",
    "- 未来工作方向：",
    "  • 进一步扩展模型规模和训练数据",
    "  • 优化奖励机制，提高推理的准确性和效率",
    "  • 探索多模态推理能力的增强",
    "  • 研究模型在更复杂任务上的表现"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    if "•" in point:
        p.level = 1
    else:
        p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)

# 致谢页
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "谢谢聆听！"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 添加联系方式
left = Inches(3)
top = Inches(4)
width = Inches(4)
height = Inches(1)

shape = slide.shapes.add_textbox(left, top, width, height)
tf = shape.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "论文链接：https://arxiv.org/pdf/2501.12948"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(102, 102, 102)
p.alignment = PP_ALIGN.CENTER

# 保存PPT
ppt_path = "DeepSeek-R1_Research_Presentation.pptx"
prs.save(ppt_path)
print(f"PPT已成功保存到: {ppt_path}")
